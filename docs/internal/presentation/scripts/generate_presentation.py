#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Erzeugt eine kleine AbschlussprÃ¤sentation fÃ¼r MinAn V1 als PPTX und PDF."""

from __future__ import annotations

import os
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

os.environ.setdefault("QT_QPA_PLATFORM", "windows")

PROJECT_DIR = Path(__file__).resolve().parents[4]
SRC_DIR = PROJECT_DIR / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from PySide6.QtCore import QTimer
from PySide6.QtWidgets import QApplication

from minan_v1.domain.session_state import SessionState
from minan_v1.services.import_service import load_csv
from minan_v1.services.profile_service import create_profile
from minan_v1.services.quality_service import compute_quality_report
from minan_v1.services.summary_service import generate_summary
from minan_v1.ui.main_window import MainWindow


PRESENTATION_DIR = PROJECT_DIR / "docs" / "internal" / "presentation"
ASSET_DIR = PRESENTATION_DIR / "generated_assets"
PPTX_PATH = PRESENTATION_DIR / "MinAn_1_4_Abschlusspraesentation.pptx"
PDF_PATH = PRESENTATION_DIR / "MinAn_1_4_Abschlusspraesentation.pdf"
SAMPLE_CSV = PROJECT_DIR / "assets" / "sample_data" / "beispiel_komma.csv"

SLIDE_W = 1600
SLIDE_H = 900
HEADER_H = 112
MARGIN_X = 86
MARGIN_Y = 62

BG = (246, 248, 251)
PANEL = (255, 255, 255)
PRIMARY = (32, 79, 118)
ACCENT = (84, 124, 82)
TEXT = (38, 46, 56)
MUTED = (96, 106, 117)
BORDER = (216, 223, 231)
WARN = (209, 92, 39)

FONT_REG = PROJECT_DIR / "dist" / "MinAn_1_4" / "_internal" / "matplotlib" / "mpl-data" / "fonts" / "ttf" / "DejaVuSans.ttf"
FONT_BOLD = PROJECT_DIR / "dist" / "MinAn_1_4" / "_internal" / "matplotlib" / "mpl-data" / "fonts" / "ttf" / "DejaVuSans-Bold.ttf"


@dataclass
class SlideSpec:
    title: str
    bullets: list[str]
    screenshot: Path | None = None
    footer: str | None = None


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_BOLD if bold else FONT_REG
    return ImageFont.truetype(str(path), size=size)


def ensure_dirs() -> None:
    PRESENTATION_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(exist_ok=True)


def qcolor_to_rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        trial = f"{current} {word}"
        if draw.textbbox((0, 0), trial, font=font)[2] <= max_width:
            current = trial
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def build_dataset_context() -> dict[str, object]:
    df, import_result = load_csv(SAMPLE_CSV)
    if not import_result.success:
        raise RuntimeError(f"Beispiel-CSV konnte nicht geladen werden: {import_result.error}")

    profile = create_profile(df)
    profile.source_path = SAMPLE_CSV
    profile.file_size_bytes = SAMPLE_CSV.stat().st_size
    profile.encoding_detected = import_result.encoding
    profile.separator_detected = import_result.separator
    quality = compute_quality_report(df, profile)
    summary = generate_summary(profile, quality)

    return {
        "df": df,
        "import_result": import_result,
        "profile": profile,
        "quality": quality,
        "summary": summary,
    }


def capture_ui_screenshots(context: dict[str, object]) -> dict[str, Path]:
    app = QApplication.instance() or QApplication([])
    session = SessionState()
    window = MainWindow(session)
    window.resize(1440, 900)
    window.show()
    app.processEvents()

    df = context["df"]
    import_result = context["import_result"]
    profile = context["profile"]
    quality = context["quality"]
    summary = context["summary"]

    session.load(df, SAMPLE_CSV)
    session.import_result = import_result
    session.profile = profile
    session.quality_report = quality
    session.summary = summary

    window._overview_panel.update_overview(import_result, profile, quality, summary)
    window._table_panel.set_dataframe(session.working_df, session.hidden_columns)
    window._metrics_panel.update_metrics(profile)
    window._charts_panel.set_data(session.working_df, profile)
    window._edit_panel.update_columns(session.working_df.columns.tolist())
    window._export_panel.update_export_info()
    window.statusBar().showMessage(
        f"{import_result.file_name} geladen - {profile.row_count} Zeilen, {profile.column_count} Spalten"
    )
    app.processEvents()

    shots: dict[str, Path] = {}
    tab_widget = window._tabs

    def save_shot(name: str, tab_index: int, chart_index: int | None = None) -> None:
        tab_widget.setCurrentIndex(tab_index)
        app.processEvents()
        if chart_index is not None:
            window._charts_panel._type_combo.setCurrentIndex(chart_index)
            app.processEvents()
        QTimer.singleShot(50, app.quit)
        app.exec()
        path = ASSET_DIR / f"{name}.png"
        window.grab().save(str(path))
        shots[name] = path

    save_shot("overview", 0)
    save_shot("metrics", 2)
    save_shot("charts", 3, 1)

    window.close()
    app.processEvents()
    return shots


def build_slide_specs(test_count: int, screenshots: dict[str, Path]) -> list[SlideSpec]:
    return [
        SlideSpec(
            title="MinAn V1",
            bullets=[
                "Portables Windows-Tool zur lokalen CSV-Schnellanalyse",
                "Version 1.0.0",
                "AbschlussprÃ¤sentation auf Basis des realen Projektstands",
            ],
            footer="Quellen: README, Architekturdoku, Analyseumfang, Release- und Testdokumentation",
        ),
        SlideSpec(
            title="Problem und Ziel",
            bullets=[
                "Ausgangslage: CSV-Dateien sollen schnell lesbar und lokal analysierbar sein.",
                "Zielbild von MinAn V1: lokal, schnell, verstÃ¤ndlich und ohne Pflichtinstallation.",
                "Zentrale Schutzregel: Originaldatei bleibt unverÃ¤ndert, gearbeitet wird auf einer Arbeitskopie.",
                "Ergebnis: Analyse und Export als neue CSV-Datei im selben Nutzungskontext.",
            ],
        ),
        SlideSpec(
            title="Funktionsumfang V1",
            bullets=[
                "CSV laden mit automatischer Erkennung von Encoding und Trennzeichen.",
                "Ãœberblick mit Dateiinfo, Strukturprofil, DatenqualitÃ¤t und Kurzzusammenfassung.",
                "Tabellenvorschau, Standardkennzahlen und Standarddiagramme.",
                "Bearbeiten auf Arbeitskopie: umbenennen, lÃ¶schen, ausblenden, filtern, sortieren, markieren.",
                "Export der Arbeitskopie als neue CSV-Datei ohne Ãœberschreiben der Quelle.",
            ],
        ),
        SlideSpec(
            title="Produktlogik und Architektur",
            bullets=[
                "Datenfluss: Originaldatei -> Import -> Original-DF -> Arbeitskopie -> Analyse/Bearbeitung -> Export-CSV.",
                "SessionState trennt Original und Arbeitskopie strikt voneinander.",
                "Services kapseln Import, Profilierung, QualitÃ¤tsprÃ¼fung, Transformation und Export.",
                "UI ist von der Fachlogik getrennt; Pfade bleiben relativ und portabel.",
            ],
            screenshot=screenshots.get("overview"),
            footer="Beispielansicht: Ãœberblick-Tab mit geladener Beispiel-CSV aus assets/sample_data",
        ),
        SlideSpec(
            title="OberflÃ¤che",
            bullets=[
                "Tabs: Ãœberblick, Tabelle, Kennzahlen, Diagramme, Bearbeiten, Export.",
                "Der Ãœberblick bÃ¼ndelt Strukturprofil, QualitÃ¤tsbefunde sowie Warnungen und Hinweise.",
                "Kennzahlen und Diagramme werden passend zum Datentyp bereitgestellt.",
            ],
            screenshot=screenshots.get("metrics"),
            footer="Echter Screenshot aus MinAn V1 mit der im Projekt enthaltenen Beispiel-CSV",
        ),
        SlideSpec(
            title="QualitÃ¤tssicherung",
            bullets=[
                f"Aktueller Testlauf im Projektordner: {test_count} Tests bestanden, 0 fehlgeschlagen.",
                "Testabdeckung laut Testplan: Import, Profilierung, QualitÃ¤t, Diagramme, Transformation, Export, Session-State.",
                "Smoke-Check und Produkt-Check sind in docs/RELEASE_STATUS.md als erfolgreich dokumentiert.",
                "Wesentliche Schutzregel wurde validiert: Originaldatei und original_df bleiben unverÃ¤ndert.",
            ],
            screenshot=screenshots.get("charts"),
            footer="Hinweis: docs/RELEASE_STATUS.md nennt 75 Tests; der aktuelle lokale Lauf ergab 83 Tests.",
        ),
        SlideSpec(
            title="Release, PortabilitÃ¤t und Fazit",
            bullets=[
                "Release-Form: One-Folder-Paket in dist/MinAn_1_4/ mit MinAn.exe und _internal/.",
                "Portable Nutzung ohne Installer, ohne Registry-Eintrag und ohne Schreibzugriff auf Systemverzeichnisse.",
                "V1-Ziele sind erreicht: lokale CSV-Schnellanalyse, verstÃ¤ndliche OberflÃ¤che, sichere Arbeitskopie, Export als neue Datei.",
                "Reale Restpunkte laut Doku: kein App-Icon, keine EXE-Versionsinfo, keine automatische Bereinigung, keine Drilldowns.",
            ],
            footer="Stand: Build- und Releasestatus aus den Projektdateien, keine externen Quellen",
        ),
    ]


def draw_slide(spec: SlideSpec, index: int, total: int) -> Path:
    image = Image.new("RGB", (SLIDE_W, SLIDE_H), BG)
    draw = ImageDraw.Draw(image)
    title_font = load_font(40, bold=True)
    body_font = load_font(25)
    small_font = load_font(18)
    label_font = load_font(18, bold=True)

    draw.rounded_rectangle((36, 28, SLIDE_W - 36, SLIDE_H - 28), radius=28, fill=PANEL, outline=BORDER, width=2)
    draw.rectangle((36, 28, SLIDE_W - 36, 28 + HEADER_H), fill=PRIMARY)
    draw.text((MARGIN_X, 58), spec.title, font=title_font, fill=(255, 255, 255))
    slide_no = f"{index}/{total}"
    slide_no_box = draw.textbbox((0, 0), slide_no, font=small_font)
    draw.text((SLIDE_W - MARGIN_X - (slide_no_box[2] - slide_no_box[0]), 66), slide_no, font=small_font, fill=(235, 241, 247))

    body_top = 176
    left_x = MARGIN_X
    left_w = 650 if spec.screenshot else 1320
    bullet_y = body_top
    bullet_gap = 18
    bullet_color = ACCENT

    for bullet in spec.bullets:
        lines = wrap_text(draw, bullet, body_font, left_w - 48)
        draw.ellipse((left_x, bullet_y + 12, left_x + 14, bullet_y + 26), fill=bullet_color)
        line_y = bullet_y
        for idx, line in enumerate(lines):
            text_x = left_x + 32 if idx == 0 else left_x + 32
            draw.text((text_x, line_y), line, font=body_font, fill=TEXT)
            line_y += 36
        bullet_y = line_y + bullet_gap

    if spec.screenshot and spec.screenshot.exists():
        shot = Image.open(spec.screenshot).convert("RGB")
        shot.thumbnail((700, 500))
        panel_x = 820
        panel_y = 190
        panel_w = 670
        panel_h = 470
        draw.rounded_rectangle((panel_x, panel_y, panel_x + panel_w, panel_y + panel_h), radius=24, fill=(248, 250, 252), outline=BORDER, width=2)
        shot_x = panel_x + (panel_w - shot.width) // 2
        shot_y = panel_y + (panel_h - shot.height) // 2
        image.paste(shot, (shot_x, shot_y))
        draw.text((panel_x, panel_y - 34), "Realer Anwendungsscreenshot", font=label_font, fill=MUTED)

    if spec.footer:
        draw.line((MARGIN_X, SLIDE_H - 108, SLIDE_W - MARGIN_X, SLIDE_H - 108), fill=BORDER, width=2)
        footer_fill = WARN if "75 Tests" in spec.footer else MUTED
        footer_lines = wrap_text(draw, spec.footer, small_font, SLIDE_W - (2 * MARGIN_X))
        fy = SLIDE_H - 92
        for line in footer_lines:
            draw.text((MARGIN_X, fy), line, font=small_font, fill=footer_fill)
            fy += 24

    out_path = ASSET_DIR / f"slide_{index:02d}.png"
    image.save(out_path)
    return out_path


def create_pdf(slide_images: Iterable[Path]) -> None:
    with PdfPages(PDF_PATH) as pdf:
        for path in slide_images:
            img = Image.open(path).convert("RGB")
            fig = plt.figure(figsize=(13.333, 7.5), dpi=120)
            ax = fig.add_axes([0, 0, 1, 1])
            ax.imshow(img)
            ax.axis("off")
            pdf.savefig(fig, dpi=120, bbox_inches="tight", pad_inches=0)
            plt.close(fig)


def add_textbox(slide, left, top, width, height, text, font_size, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = qcolor_to_rgb(color)
    p.alignment = align
    return box


def create_pptx(specs: list[SlideSpec]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, spec in enumerate(specs, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = qcolor_to_rgb(BG)

        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.25), Inches(12.73), Inches(7.0))
        card.fill.solid()
        card.fill.fore_color.rgb = qcolor_to_rgb(PANEL)
        card.line.color.rgb = qcolor_to_rgb(BORDER)
        card.line.width = Pt(1.25)

        header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.3), Inches(0.25), Inches(12.73), Inches(0.95))
        header.fill.solid()
        header.fill.fore_color.rgb = qcolor_to_rgb(PRIMARY)
        header.line.fill.background()

        add_textbox(slide, Inches(0.72), Inches(0.48), Inches(9.5), Inches(0.38), spec.title, 24, bold=True, color=(255, 255, 255))
        add_textbox(slide, Inches(11.7), Inches(0.5), Inches(0.8), Inches(0.3), f"{idx}/{len(specs)}", 11, color=(233, 239, 245), align=PP_ALIGN.RIGHT)

        body_left = Inches(0.72)
        body_top = Inches(1.45)
        body_width = Inches(5.5 if spec.screenshot else 11.1)

        bullet_box = slide.shapes.add_textbox(body_left, body_top, body_width, Inches(4.8))
        bullet_frame = bullet_box.text_frame
        bullet_frame.word_wrap = True
        bullet_frame.margin_left = Pt(0)
        bullet_frame.margin_right = Pt(0)
        for bullet_idx, bullet in enumerate(spec.bullets):
            p = bullet_frame.paragraphs[0] if bullet_idx == 0 else bullet_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = qcolor_to_rgb(TEXT)
            p.bullet = True
            p.space_after = Pt(10)

        if spec.screenshot and spec.screenshot.exists():
            frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.85), Inches(1.58), Inches(5.45), Inches(3.9))
            frame.fill.solid()
            frame.fill.fore_color.rgb = qcolor_to_rgb((248, 250, 252))
            frame.line.color.rgb = qcolor_to_rgb(BORDER)
            frame.line.width = Pt(1.0)
            add_textbox(slide, Inches(6.95), Inches(1.28), Inches(4.0), Inches(0.25), "Realer Anwendungsscreenshot", 10, bold=True, color=MUTED)
            slide.shapes.add_picture(str(spec.screenshot), Inches(7.0), Inches(1.82), width=Inches(5.1))

        if spec.footer:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(6.45), Inches(11.1), Inches(0.02))
            line.fill.solid()
            line.fill.fore_color.rgb = qcolor_to_rgb(BORDER)
            line.line.fill.background()
            color = WARN if "75 Tests" in spec.footer else MUTED
            add_textbox(slide, Inches(0.72), Inches(6.52), Inches(11.1), Inches(0.42), spec.footer, 10, color=color)

    prs.save(PPTX_PATH)


def run_pytest_and_get_count() -> int:
    import subprocess

    result = subprocess.run(
        [sys.executable, "-m", "pytest", "-q"],
        cwd=PROJECT_DIR,
        capture_output=True,
        text=True,
        check=True,
    )
    for line in result.stdout.splitlines():
        if "passed" in line and "failed" not in line:
            token = line.strip().split()[0]
            if token.isdigit():
                return int(token)
    raise RuntimeError("Konnte Testanzahl aus pytest-Ausgabe nicht ermitteln.")


def main() -> None:
    ensure_dirs()
    context = build_dataset_context()
    screenshots = capture_ui_screenshots(context)
    test_count = run_pytest_and_get_count()
    specs = build_slide_specs(test_count, screenshots)
    slide_images = [draw_slide(spec, idx, len(specs)) for idx, spec in enumerate(specs, start=1)]
    create_pdf(slide_images)
    create_pptx(specs)
    print(f"[OK] PDF erstellt: {PDF_PATH}")
    print(f"[OK] PPTX erstellt: {PPTX_PATH}")
    print(f"[OK] Screenshots erstellt: {', '.join(str(path) for path in screenshots.values())}")


if __name__ == "__main__":
    main()

