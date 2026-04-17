#!/usr/bin/env python3
"""
Erstelle professionelle AbschlussprÃ¤sentation fÃ¼r MinAn 1.4
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Farben
COLOR_DARK_BG = RGBColor(25, 35, 40)
COLOR_DARK_PANEL = RGBColor(35, 50, 60)
COLOR_GREEN_ACCENT = RGBColor(76, 175, 80)  # GrÃ¼n
COLOR_LIGHT_TEXT = RGBColor(230, 235, 240)
COLOR_GRAY_TEXT = RGBColor(150, 160, 170)

PROJECT_ROOT = Path(__file__).resolve().parents[4]
ASSETS_DIR = PROJECT_ROOT / "assets"
ICON_PATH = ASSETS_DIR / "icons" / "minan_v1.png"
RELEASE_SCREENSHOT = PROJECT_ROOT / "docs" / "release_assets" / "MinAn_1_4_Release_Screenshot.png"

# PowerPoint-PrÃ¤sentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    """Titelfolie mit dunkel/grÃ¼n Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dunkler Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BG

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_GREEN_ACCENT
        paragraph.alignment = PP_ALIGN.CENTER

    # Untertitel
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        for paragraph in subtitle_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = COLOR_LIGHT_TEXT
            paragraph.alignment = PP_ALIGN.CENTER

    # Icon (wenn vorhanden)
    if ICON_PATH.exists():
        try:
            slide.shapes.add_picture(str(ICON_PATH), Inches(4), Inches(0.5), width=Inches(2))
        except:
            pass

    return slide

def add_content_slide(title, content_points=None, image_path=None):
    """Inhaltsfolie mit Titel und Bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BG

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_GREEN_ACCENT

    # GrÃ¼ne Trennlinie
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = COLOR_GREEN_ACCENT
    line.line.width = Pt(2)

    # Inhalte
    if content_points:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for point in content_points:
            if isinstance(point, tuple):
                text, is_bold = point
            else:
                text = point
                is_bold = False

            p = text_frame.add_paragraph()
            p.text = text
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_LIGHT_TEXT
            p.font.bold = is_bold
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    # Bild rechts (wenn vorhanden)
    if image_path and Path(image_path).exists():
        try:
            slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4))
        except:
            pass

    return slide

def add_two_column_slide(title, left_content, right_content):
    """Zwei-Spalten-Layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BG

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_GREEN_ACCENT

    # GrÃ¼ne Trennlinie
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = COLOR_GREEN_ACCENT
    line.line.width = Pt(2)

    # Linke Spalte
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, point in enumerate(left_content):
        if isinstance(point, tuple):
            text, is_bold = point
        else:
            text = point
            is_bold = False

        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = is_bold
        p.font.color.rgb = COLOR_LIGHT_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Rechte Spalte
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, point in enumerate(right_content):
        if isinstance(point, tuple):
            text, is_bold = point
        else:
            text = point
            is_bold = False

        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = is_bold
        p.font.color.rgb = COLOR_LIGHT_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return slide

# ============================================================================
# FOLIEN ERSTELLEN
# ============================================================================

# 1. TITELFOLIE
add_title_slide(
    "MinAn 1.4",
    "CSV-Schnellanalyse\n\nLokales, portables Windows-Tool zur schnellen Analyse und Bearbeitung von CSV-Dateien"
)

# 2. AUSGANGSLAGE / PROBLEM
add_content_slide(
    "Ausgangslage",
    [
        "â€¢ Manuelle CSV-Erstanalyse ist zeitintensiv",
        "â€¢ Unstrukturierter Ãœberblick: Struktur, QualitÃ¤t, AusreiÃŸer",
        "â€¢ Keine lokale, portable LÃ¶sung ohne Installation",
        "â€¢ Cloud-AbhÃ¤ngigkeit oft nicht gewÃ¼nscht",
        "â€¢ Originaldatei muss geschÃ¼tzt sein",
        "â†’ LÃ¶sung: Schneller, lokaler, offline Ãœberblick"
    ]
)

# 3. PRODUKTZIEL
add_content_slide(
    "Produktziel",
    [
        "âœ“ Schneller Ãœberblick ohne Setup",
        "âœ“ VollstÃ¤ndig portabel und lokal",
        "âœ“ Originaldatei geschÃ¼tzt",
        "âœ“ Arbeitsfluss fÃ¼r Analysten: Laden â†’ PrÃ¼fen â†’ Bearbeiten â†’ Export",
        "âœ“ Fokus auf Fachlichkeit, nicht KomplexitÃ¤t",
        "âœ“ Windows-nativer Desktop-Workflow"
    ]
)

# 4. FUNKTIONSUMFANG
add_content_slide(
    "Funktionsumfang",
    [
        "CSV-Verwaltung:",
        "  â€¢ CSV laden mit automatischer Encoding/Separator-Erkennung",
        "  â€¢ Beispieldatei (Schnellstart)",
        "",
        "Analyse & Vorschau:",
        "  â€¢ Strukturprofil und DatenqualitÃ¤t",
        "  â€¢ Kennzahlen und Diagramme",
        "  â€¢ Tabellenansicht mit Scrolling",
        "",
        "Arbeitstools:",
        "  â€¢ Mehrfach-Filter, Schnellansichten",
        "  â€¢ Spalten-Bearbeitung (Umbenennen, LÃ¶schen, Typ-Overrides)",
        "  â€¢ CSV-Export & HTML-Bericht (aktive Sicht)"
    ]
)

# 5. PRODUKTLOGIK
add_content_slide(
    "Produktlogik & Schutzkonzept",
    [
        "Originaldatei â†’ Arbeitskopie",
        "  â€¢ Originaldatei wird beim Laden gespeichert",
        "  â€¢ Alle Ã„nderungen nur auf Arbeitskopie",
        "",
        "Aktive Sicht",
        "  â€¢ Filter und Transformationen definieren aktive Sicht",
        "  â€¢ Kennzahlen, Diagramme, Export beziehen sich auf aktive Sicht",
        "",
        "Export & Bericht",
        "  â€¢ CSV-Export speichert aktive Sicht als neue Datei",
        "  â€¢ HTML-Bericht dokumentiert aktive Sicht",
        "",
        "Schutz",
        "  â€¢ Originaldatei wird nie verÃ¤ndert oder Ã¼berschrieben"
    ]
)

# 6. OBERFLÃ„CHE / PRODUKTANSICHT
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_DARK_BG

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "ProduktoberflÃ¤che"
for paragraph in title_frame.paragraphs:
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.color.rgb = COLOR_GREEN_ACCENT

line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
line.line.color.rgb = COLOR_GREEN_ACCENT
line.line.width = Pt(2)

# Screenshot
if RELEASE_SCREENSHOT.exists():
    try:
        slide.shapes.add_picture(str(RELEASE_SCREENSHOT), Inches(0.5), Inches(1.3), width=Inches(9))
    except Exception as e:
        # Fallback: Text
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = text_box.text_frame
        tf.text = "â€¢ Toolbar: CSV Ã¶ffnen, Bericht exportieren, Info/Schnellstart\nâ€¢ Tabs: Ãœberblick, Kennzahlen, Diagramme, Tabelle, Bearbeiten, Export\nâ€¢ Dunkles ModernDesign mit grÃ¼nen Akzenten\nâ€¢ Responsive Layouts fÃ¼r groÃŸe und kleine Dateien"
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = COLOR_LIGHT_TEXT

# 7. TECHNIKBLOCK - Stack
add_content_slide(
    "Technik-Stack",
    [
        ("Programmiersprache", True),
        "  Python 3.10+",
        "",
        ("UI-Framework", True),
        "  PySide6 (Qt6 fÃ¼r Python)",
        "",
        ("Datenverwaltung", True),
        "  pandas (DataFrames, Analyse)",
        "  numpy (numerische Berechnung)",
        "",
        ("Visualisierung", True),
        "  matplotlib (Diagramme, Charts)",
        "",
        ("QualitÃ¤tssicherung", True),
        "  pytest (automatisierte Tests)"
    ]
)

# 8. TECHNIKBLOCK - Build & Release
add_content_slide(
    "Build & Release",
    [
        ("Packaging", True),
        "  PyInstaller One-Folder-Release",
        "  â†’ MinAn_1_4/ Ordner (portable)",
        "",
        ("Release-Struktur", True),
        "  â”œ MinAn.exe (Hauptanwendung)",
        "  â”œ _internal/ (Python Runtime, Libraries)",
        "  â”œ output/ (Berichte und CSV-Exporte)",
        "  â”” _internal/sample_data/ (Beispieldatei)",
        "",
        ("Versionierung", True),
        "  Windows Version Info: 1.4.0.0",
        "  Productname: MinAn 1.4 - CSV-Schnellanalyse"
    ]
)

# 9. ARCHITEKTUR
add_two_column_slide(
    "Projektarchitektur",
    [
        ("Schichtmodell", True),
        "",
        "UI-Layer",
        "  main_window",
        "  dialogs, widgets",
        "  Qt-Models",
        "",
        "Services",
        "  import, profile",
        "  quality, chart",
        "  transform, export",
        "  report"
    ],
    [
        ("SessionState & Domain", True),
        "",
        "SessionState",
        "  original_df",
        "  working_df",
        "  filters & views",
        "",
        "Domain Models",
        "  ColumnProfile",
        "  DataQuality",
        "  enums.py"
    ]
)

# 10. QUALITÃ„TSSICHERUNG
add_content_slide(
    "QualitÃ¤tssicherung",
    [
        ("Automatisierte Tests", True),
        "  44+ Pytest-Tests abdecken:",
        "  â€¢ CSV-Import mit verschiedenen Encodings",
        "  â€¢ Profilierung (Zeilen, Spalten, Typen)",
        "  â€¢ DatenqualitÃ¤t (Missing, Dubletten)",
        "  â€¢ Transformationen (Filter, Markierungen)",
        "  â€¢ Export und Bericht-Generierung",
        "  â€¢ Originalschutz und Session-State",
        "",
        ("Manuelle Smoke-Tests", True),
        "  âœ“ App startet und lÃ¤dt Dateien",
        "  âœ“ Alle Tabs funktionieren",
        "  âœ“ Export und Bericht-Generierung"
    ]
)

# 11. RELEASE & PORTABILITÃ„T
add_content_slide(
    "Release & PortabilitÃ¤t",
    [
        ("One-Folder-Prinzip", True),
        "  dist/MinAn_1_4/ ist vollstÃ¤ndig portabel",
        "  Keine Systeminstallation erforderlich",
        "",
        ("Ausgabeordner", True),
        "  output/reports/ â†’ HTML-Berichte",
        "  output/csv/ â†’ CSV-Exporte",
        "",
        ("Beispieldatei", True),
        "  test_csv_deutsch_200x15.csv",
        "  200 Zeilen, 15 Spalten, deutsch",
        "  â†’ Schnellstart und Demo",
        "",
        ("Schnellstart", True),
        "  Info-Dialog mit Beispieldatei-Button"
    ]
)

# 12. PRODUKTREIFE & FEINSCHLIFF
add_content_slide(
    "Produktreife - Feinschliff",
    [
        ("UI/UX Politur", True),
        "  âœ“ Dark Mode mit grÃ¼nen Akzenten",
        "  âœ“ Responsive Layouts",
        "  âœ“ Klare, konsistente Toolbar",
        "  âœ“ Info-Dialog und Schnellstart",
        "",
        ("Produktmetadaten", True),
        "  âœ“ Produktname: MinAn 1.4 - CSV-Schnellanalyse",
        "  âœ“ Company: MinAn Software",
        "  âœ“ Icon: ModernGreen Analyse-Symbol",
        "",
        ("Fertigstellung 1.4", True),
        "  âœ“ Alle Tests grÃ¼n",
        "  âœ“ Release-Build erfolgreich",
        "  âœ“ Portable EXE lauffÃ¤hig"
    ]
)

# 13. FAZIT
add_content_slide(
    "Fazit",
    [
        ("Was ist MinAn 1.4?", True),
        "Ein reifes, produktives Analysewerkzeug fÃ¼r CSV-Daten",
        "",
        ("KernstÃ¤rken", True),
        "  âœ“ Lokal, portabel, sicher",
        "  âœ“ Schneller Ãœberblick ohne Setup",
        "  âœ“ Originalschutz durch Arbeitskopien-Prinzip",
        "  âœ“ Professioneller Export und Bericht",
        "  âœ“ Solide Test-Basis und Release-QualitÃ¤t",
        "",
        ("Projektstand", True),
        "  âœ“ Version 1.4 ist Abschluss von Blocks 1â€“2",
        "  âœ“ Produktionsreifer Stand",
        "  âœ“ Bereit fÃ¼r Einsatz und Weiterentwicklung"
    ]
)

# SPEICHERN
output_path = PROJECT_ROOT / "docs" / "internal" / "presentation" / "MinAn_1_4_Abschlusspraesentation.pptx"
output_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(output_path))

print(f"[OK] PrÃ¤sentation erstellt: {output_path}")
print(f"[OK] Folien: {len(prs.slides)}")

