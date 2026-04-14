#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Erzeugt die aktuelle Abschlusspräsentation aus dem Stand MinAn 1.4."""

from __future__ import annotations

import os
import re
import shutil
import subprocess
import sys
from dataclasses import dataclass
from pathlib import Path

os.environ.setdefault("QT_QPA_PLATFORM", "windows")

PROJECT_DIR = Path(__file__).resolve().parents[2]
SRC_DIR = PROJECT_DIR / "src"
if str(SRC_DIR) not in sys.path:
    sys.path.insert(0, str(SRC_DIR))

from matplotlib.backends.backend_pdf import PdfPages
import matplotlib.pyplot as plt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PySide6.QtWidgets import QApplication

from minan_v1.domain.session_state import SessionState
from minan_v1.resources import sample_data_path
from minan_v1.ui.main_window import MainWindow


PRESENTATION_DIR = PROJECT_DIR / "presentation"
ASSET_DIR = PRESENTATION_DIR / "generated_assets"
SLIDE_DIR = ASSET_DIR / "current_slides"

OUTPUTS = {
    "pdf_main": PRESENTATION_DIR / "MinAn_1_4_Abschlusspraesentation.pdf",
    "pptx_main": PRESENTATION_DIR / "MinAn_1_4_Abschlusspraesentation.pptx",
    "pdf_current": PRESENTATION_DIR / "MinAn_1_4_Abschlusspraesentation.pdf",
    "pptx_current": PRESENTATION_DIR / "MinAn_1_4_Abschlusspraesentation.pptx",
}

FONT_REG = PROJECT_DIR / "dist" / "MinAn_1_4" / "_internal" / "matplotlib" / "mpl-data" / "fonts" / "ttf" / "DejaVuSans.ttf"
FONT_BOLD = PROJECT_DIR / "dist" / "MinAn_1_4" / "_internal" / "matplotlib" / "mpl-data" / "fonts" / "ttf" / "DejaVuSans-Bold.ttf"

W = 1600
H = 900

BG = (245, 247, 245)
PANEL = (255, 255, 255)
PRIMARY = (46, 88, 82)
ACCENT = (117, 155, 119)
TEXT = (36, 46, 45)
MUTED = (96, 110, 107)
BORDER = (212, 220, 216)
WARN = (176, 110, 45)


@dataclass
class Slide:
    title: str
    bullets: list[str]
    images: list[tuple[Path, str]]
    footer: str | None = None


def ensure_dirs() -> None:
    PRESENTATION_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(exist_ok=True)
    SLIDE_DIR.mkdir(exist_ok=True)


def font(size: int, bold: bool = False):
    path = FONT_BOLD if bold else FONT_REG
    return ImageFont.truetype(str(path), size)


def wrap(draw: ImageDraw.ImageDraw, text: str, fnt, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]
    lines = []
    current = words[0]
    for word in words[1:]:
        probe = f"{current} {word}"
        if draw.textbbox((0, 0), probe, font=fnt)[2] <= max_width:
            current = probe
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def parse_pytest_count() -> int:
    result = subprocess.run(
        [sys.executable, "-m", "pytest", "-q"],
        cwd=PROJECT_DIR,
        capture_output=True,
        text=True,
        check=True,
    )
    match = re.search(r"(\d+)\s+passed", result.stdout)
    if not match:
        raise RuntimeError("Konnte Testanzahl nicht aus pytest-Ausgabe lesen.")
    return int(match.group(1))


def exe_size_mb() -> float:
    exe = PROJECT_DIR / "dist" / "MinAn_1_4" / "MinAn.exe"
    return round(exe.stat().st_size / (1024 * 1024), 1)


def capture_screenshots() -> dict[str, Path]:
    app = QApplication.instance() or QApplication([])
    session = SessionState()
    window = MainWindow(session)
    window.resize(1440, 920)
    window.show()
    app.processEvents()

    sample = sample_data_path()
    if not sample.exists():
        raise FileNotFoundError(f"Beispieldatei fehlt: {sample}")
    if not window.load_csv_path(sample):
        raise RuntimeError("Beispieldatei konnte nicht geladen werden.")
    app.processEvents()

    targets = {
        "overview": 0,
        "metrics": 2,
        "charts": 3,
        "edit": 4,
    }
    shots: dict[str, Path] = {}
    for name, tab_index in targets.items():
        window._tabs.setCurrentIndex(tab_index)
        app.processEvents()
        if name == "charts":
            window._charts_panel._type_combo.setCurrentIndex(1)
            app.processEvents()
        path = ASSET_DIR / f"{name}_current.png"
        window.grab().save(str(path))
        shots[name] = path

    window.close()
    app.quit()
    return shots


def build_slides(test_count: int, shots: dict[str, Path]) -> list[Slide]:
    return [
        Slide(
            title="MinAn 1.4",
            bullets=[
                "Aktueller Projektstand im Ordner MinAn V1.",
                "Portables Windows-Tool zur lokalen CSV-Schnellanalyse.",
                "Fokus laut Projektstand: Schnellstart, aktive Sicht, CSV-Export und HTML-Bericht.",
            ],
            images=[],
            footer="Basis: aktueller README- und Doku-Stand vom 07.04.2026",
        ),
        Slide(
            title="Problem und Ziel",
            bullets=[
                "CSV-Dateien sollen ohne Installation lokal geladen und schnell eingeordnet werden.",
                "Der erste Blick auf Struktur, Datenqualität, Kennzahlen und Diagramme soll direkt verfügbar sein.",
                "Die Originaldatei bleibt unangetastet; gearbeitet wird auf einer Arbeitskopie und aktiven Sicht.",
                "Nutzdaten sollen im portablen Release sichtbar im Hauptordner landen, nicht versteckt in Systempfaden.",
            ],
            images=[],
        ),
        Slide(
            title="Funktionsumfang 1.4",
            bullets=[
                "CSV laden mit Encoding- und Separator-Erkennung.",
                "Überblick, Tabelle, Kennzahlen und Diagramme nach dem Laden.",
                "Bearbeiten auf Arbeitskopie: Mehrfachfilter, Schnellansichten, Markierungen und Typübersteuerung.",
                "Export der aktiven Sicht als neue CSV-Datei.",
                "Lokaler HTML-Analysebericht aus derselben aktiven Sicht.",
                "Info / Schnellstart mit Button für die mitgelieferte Beispieldatei.",
            ],
            images=[],
        ),
        Slide(
            title="Produktlogik und Architektur",
            bullets=[
                "Datenfluss: CSV-Datei -> import_service -> Original-DF -> Arbeitskopie -> aktive Sicht -> CSV/HTML-Ausgabe.",
                "SessionState trennt Originaldaten, Arbeitskopie und aktive Sicht.",
                "Services kapseln Import, Profilierung, Qualität, Diagramme, Transformation, Export und Bericht.",
                "UI, Services und Domain sind laut Architektur getrennt; relevante Pfade bleiben portabel.",
            ],
            images=[(shots["edit"], "Bearbeiten-Tab mit Filtern, Schnellansichten und Markierungen")],
            footer="Realer Screenshot aus der aktuellen Anwendung mit test_csv_deutsch_200x15.csv",
        ),
        Slide(
            title="Oberfläche und Screenshots",
            bullets=[
                "Die Oberfläche bündelt Überblick, Kennzahlen und Diagramme in getrennten Tabs.",
                "Die mitgelieferte Beispieldatei umfasst 200 Zeilen und 15 Spalten.",
                "Im Überblick werden Struktur, Warnungen und Hinweise direkt nach dem Laden sichtbar.",
            ],
            images=[
                (shots["overview"], "Überblick"),
                (shots["metrics"], "Kennzahlen"),
                (shots["charts"], "Diagramme"),
            ],
            footer="Verwendet wurden neu erzeugte reale Screenshots aus dem aktuellen Projektstand.",
        ),
        Slide(
            title="Qualitätssicherung",
            bullets=[
                f"Aktueller lokaler Testlauf im Projektordner: {test_count} Tests bestanden.",
                "Testbereiche laut Testplan: Import, Profilierung, Datenqualität, Diagramme, Transformationen, Export/Bericht sowie Session-State und Pfade.",
                "Release-Status nennt produktrelevante Tests, erfolgreichen finalen Build und startende EXE.",
                "Die Schutzregel Originaldatei bleibt unverändert ist in Doku und Tests durchgängig verankert.",
            ],
            images=[],
            footer="Realer Testlauf: pytest -q im Projektordner",
        ),
        Slide(
            title="Release und Portabilität",
            bullets=[
                "One-Folder-Release: dist/MinAn_1_4/ mit MinAn.exe, _internal, output und README_Kurzstart.txt.",
                f"Startdatei: MinAn.exe mit ca. {exe_size_mb()} MB.",
                "Output-Ordner im Release: output/reports für HTML-Berichte und output/csv für CSV-Exporte.",
                "Die Beispieldatei liegt im Release unter _internal/sample_data/test_csv_deutsch_200x15.csv.",
                "Keine Pflichtinstallation und keine Nutzdatenablage in _internal außer der mitgelieferten Beispieldatei.",
            ],
            images=[],
        ),
        Slide(
            title="Fazit",
            bullets=[
                "Der aktuelle Stand ist kein reines V1-Minimum mehr, sondern dokumentiert als MinAn 1.4.",
                "Erreicht wurden lokaler Schnellstart, portable Release-Struktur, aktive Sicht, CSV-Export und HTML-Bericht.",
                "Die Architektur bleibt klar getrennt, die Originaldatei bleibt geschützt und der Release ist portable nutzbar.",
                "Offene Restpunkte wurden in den aktuellen Kerndokumenten nicht als release-blockend ausgewiesen.",
            ],
            images=[],
            footer="Hinweis: Projektordner heißt weiterhin MinAn V1, die aktuelle Produktbezeichnung in den Quellen ist MinAn 1.4.",
        ),
    ]


def paste_contained(base: Image.Image, img_path: Path, box: tuple[int, int, int, int]) -> None:
    img = Image.open(img_path).convert("RGB")
    x, y, w, h = box
    img.thumbnail((w, h))
    px = x + (w - img.width) // 2
    py = y + (h - img.height) // 2
    base.paste(img, (px, py))


def draw_slide(slide: Slide, index: int, total: int) -> Path:
    img = Image.new("RGB", (W, H), BG)
    draw = ImageDraw.Draw(img)
    title_f = font(42, True)
    body_f = font(25)
    small_f = font(18)
    cap_f = font(18, True)

    draw.rounded_rectangle((34, 28, W - 34, H - 28), radius=26, fill=PANEL, outline=BORDER, width=2)
    draw.rounded_rectangle((34, 28, W - 34, 130), radius=26, fill=PRIMARY, outline=PRIMARY)
    draw.rectangle((34, 100, W - 34, 130), fill=PRIMARY)
    draw.text((82, 56), slide.title, fill=(255, 255, 255), font=title_f)
    draw.text((W - 120, 62), f"{index}/{total}", fill=(238, 245, 241), font=small_f)

    left_x = 82
    bullet_y = 170
    text_w = 690 if slide.images else 1370

    for bullet in slide.bullets:
        lines = wrap(draw, bullet, body_f, text_w - 40)
        draw.ellipse((left_x, bullet_y + 12, left_x + 14, bullet_y + 26), fill=ACCENT)
        line_y = bullet_y
        for line in lines:
            draw.text((left_x + 30, line_y), line, fill=TEXT, font=body_f)
            line_y += 36
        bullet_y = line_y + 14

    if slide.images:
        if len(slide.images) == 1:
            path, caption = slide.images[0]
            box = (820, 176, 690, 520)
            draw.rounded_rectangle((box[0] - 10, box[1] - 34, box[0] + box[2] + 10, box[1] + box[3] + 10), radius=20, fill=(250, 251, 250), outline=BORDER, width=2)
            draw.text((box[0], box[1] - 26), caption, fill=MUTED, font=cap_f)
            paste_contained(img, path, box)
        else:
            thumb_boxes = [
                (770, 210, 360, 220),
                (1140, 210, 360, 220),
                (955, 470, 360, 220),
            ]
            for (path, caption), box in zip(slide.images, thumb_boxes):
                x, y, w, h = box
                draw.rounded_rectangle((x - 8, y - 28, x + w + 8, y + h + 8), radius=16, fill=(250, 251, 250), outline=BORDER, width=2)
                draw.text((x, y - 22), caption, fill=MUTED, font=cap_f)
                paste_contained(img, path, box)

    if slide.footer:
        draw.line((82, H - 104, W - 82, H - 104), fill=BORDER, width=2)
        footer_color = WARN if "pytest -q" in slide.footer else MUTED
        fy = H - 88
        for line in wrap(draw, slide.footer, small_f, W - 164):
            draw.text((82, fy), line, fill=footer_color, font=small_f)
            fy += 24

    out = SLIDE_DIR / f"slide_{index:02d}.png"
    img.save(out)
    return out


def create_pdf(slide_paths: list[Path], target: Path) -> None:
    with PdfPages(target) as pdf:
        for path in slide_paths:
            page = Image.open(path).convert("RGB")
            fig = plt.figure(figsize=(13.333, 7.5), dpi=120)
            ax = fig.add_axes([0, 0, 1, 1])
            ax.imshow(page)
            ax.axis("off")
            pdf.savefig(fig, dpi=120, bbox_inches="tight", pad_inches=0)
            plt.close(fig)


def rgb(rgb_tuple: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*rgb_tuple)


def add_textbox(slide, left, top, width, height, text, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align
    return box


def create_pptx(slides: list[Slide], target: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for index, data in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(BG)

        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.28), Inches(0.24), Inches(12.77), Inches(7.0))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb(PANEL)
        card.line.color.rgb = rgb(BORDER)

        header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.28), Inches(0.24), Inches(12.77), Inches(0.84))
        header.fill.solid()
        header.fill.fore_color.rgb = rgb(PRIMARY)
        header.line.fill.background()
        header_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.28), Inches(0.84), Inches(12.77), Inches(0.24))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = rgb(PRIMARY)
        header_bar.line.fill.background()

        add_textbox(slide, Inches(0.68), Inches(0.47), Inches(8.5), Inches(0.4), data.title, 24, color=(255, 255, 255), bold=True)
        add_textbox(slide, Inches(11.85), Inches(0.5), Inches(0.7), Inches(0.2), f"{index}/{len(slides)}", 11, color=(237, 244, 241), align=PP_ALIGN.RIGHT)

        body_w = Inches(5.4 if data.images else 11.1)
        textbox = slide.shapes.add_textbox(Inches(0.68), Inches(1.43), body_w, Inches(4.9))
        frame = textbox.text_frame
        frame.word_wrap = True
        for bullet_index, bullet in enumerate(data.bullets):
            p = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
            p.text = bullet
            p.bullet = True
            p.font.size = Pt(18)
            p.font.color.rgb = rgb(TEXT)
            p.space_after = Pt(10)

        if data.images:
            if len(data.images) == 1:
                path, caption = data.images[0]
                frame_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.55), Inches(5.55), Inches(4.55))
                frame_shape.fill.solid()
                frame_shape.fill.fore_color.rgb = rgb((250, 251, 250))
                frame_shape.line.color.rgb = rgb(BORDER)
                add_textbox(slide, Inches(6.9), Inches(1.28), Inches(4.4), Inches(0.2), caption, 10, color=MUTED, bold=True)
                slide.shapes.add_picture(str(path), Inches(6.95), Inches(1.8), width=Inches(5.2))
            else:
                positions = [
                    (Inches(6.4), Inches(1.85), Inches(2.75), Inches(1.65)),
                    (Inches(9.3), Inches(1.85), Inches(2.75), Inches(1.65)),
                    (Inches(7.85), Inches(4.1), Inches(2.75), Inches(1.65)),
                ]
                for (path, caption), (left, top, width, height) in zip(data.images, positions):
                    frame_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left - Inches(0.05), top - Inches(0.22), width + Inches(0.1), height + Inches(0.27))
                    frame_shape.fill.solid()
                    frame_shape.fill.fore_color.rgb = rgb((250, 251, 250))
                    frame_shape.line.color.rgb = rgb(BORDER)
                    add_textbox(slide, left, top - Inches(0.18), width, Inches(0.15), caption, 9, color=MUTED, bold=True)
                    slide.shapes.add_picture(str(path), left, top, width=width)

        if data.footer:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.68), Inches(6.45), Inches(11.05), Inches(0.02))
            line.fill.solid()
            line.fill.fore_color.rgb = rgb(BORDER)
            line.line.fill.background()
            footer_color = WARN if "pytest -q" in data.footer else MUTED
            add_textbox(slide, Inches(0.68), Inches(6.53), Inches(11.05), Inches(0.34), data.footer, 10, color=footer_color)

    prs.save(target)


def duplicate_outputs() -> None:
    shutil.copy2(OUTPUTS["pdf_main"], OUTPUTS["pdf_current"])
    shutil.copy2(OUTPUTS["pptx_main"], OUTPUTS["pptx_current"])


def main() -> None:
    ensure_dirs()
    shots = capture_screenshots()
    test_count = parse_pytest_count()
    slides = build_slides(test_count, shots)
    slide_paths = [draw_slide(slide, index, len(slides)) for index, slide in enumerate(slides, start=1)]
    create_pdf(slide_paths, OUTPUTS["pdf_main"])
    create_pptx(slides, OUTPUTS["pptx_main"])
    duplicate_outputs()
    print(f"[OK] PDF erstellt: {OUTPUTS['pdf_main']}")
    print(f"[OK] PPTX erstellt: {OUTPUTS['pptx_main']}")
    print(f"[OK] Weitere Ausgabe: {OUTPUTS['pdf_current']}")
    print(f"[OK] Screenshots: {', '.join(str(path) for path in shots.values())}")


if __name__ == "__main__":
    main()
