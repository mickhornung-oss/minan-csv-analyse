#!/usr/bin/env python3
"""
Erstelle professionelle Abschlusspräsentation für MinAn 1.4
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path

# Farben
COLOR_DARK_BG = RGBColor(25, 35, 40)
COLOR_DARK_PANEL = RGBColor(35, 50, 60)
COLOR_GREEN_ACCENT = RGBColor(76, 175, 80)  # Grün
COLOR_LIGHT_TEXT = RGBColor(230, 235, 240)
COLOR_GRAY_TEXT = RGBColor(150, 160, 170)

PROJECT_ROOT = Path(__file__).resolve().parents[2]
ASSETS_DIR = PROJECT_ROOT / "assets"
ICON_PATH = ASSETS_DIR / "icons" / "minan_v1.png"
RELEASE_SCREENSHOT = PROJECT_ROOT / "docs" / "release_assets" / "MinAn_1_4_Release_Screenshot.png"

# PowerPoint-Präsentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    """Titelfolie mit dunkel/grün Design"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dunkler Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BG

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(54)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_GREEN_ACCENT
        paragraph.alignment = PP_ALIGN.CENTER

    # Untertitel
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        for paragraph in subtitle_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.color.rgb = COLOR_LIGHT_TEXT
            paragraph.alignment = PP_ALIGN.CENTER

    # Icon (wenn vorhanden)
    if ICON_PATH.exists():
        try:
            slide.shapes.add_picture(str(ICON_PATH), Inches(4), Inches(0.5), width=Inches(2))
        except:
            pass

    return slide

def add_content_slide(title, content_points=None, image_path=None):
    """Inhaltsfolie mit Titel und Bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Hintergrund
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BG

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_GREEN_ACCENT

    # Grüne Trennlinie
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = COLOR_GREEN_ACCENT
    line.line.width = Pt(2)

    # Inhalte
    if content_points:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for point in content_points:
            if isinstance(point, tuple):
                text, is_bold = point
            else:
                text = point
                is_bold = False

            p = text_frame.add_paragraph()
            p.text = text
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_LIGHT_TEXT
            p.font.bold = is_bold
            p.space_before = Pt(6)
            p.space_after = Pt(6)

    # Bild rechts (wenn vorhanden)
    if image_path and Path(image_path).exists():
        try:
            slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4))
        except:
            pass

    return slide

def add_two_column_slide(title, left_content, right_content):
    """Zwei-Spalten-Layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BG

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(40)
        paragraph.font.bold = True
        paragraph.font.color.rgb = COLOR_GREEN_ACCENT

    # Grüne Trennlinie
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = COLOR_GREEN_ACCENT
    line.line.width = Pt(2)

    # Linke Spalte
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, point in enumerate(left_content):
        if isinstance(point, tuple):
            text, is_bold = point
        else:
            text = point
            is_bold = False

        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = is_bold
        p.font.color.rgb = COLOR_LIGHT_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Rechte Spalte
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, point in enumerate(right_content):
        if isinstance(point, tuple):
            text, is_bold = point
        else:
            text = point
            is_bold = False

        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = is_bold
        p.font.color.rgb = COLOR_LIGHT_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    return slide

# ============================================================================
# FOLIEN ERSTELLEN
# ============================================================================

# 1. TITELFOLIE
add_title_slide(
    "MinAn 1.4",
    "CSV-Schnellanalyse\n\nLokales, portables Windows-Tool zur schnellen Analyse und Bearbeitung von CSV-Dateien"
)

# 2. AUSGANGSLAGE / PROBLEM
add_content_slide(
    "Ausgangslage",
    [
        "• Manuelle CSV-Erstanalyse ist zeitintensiv",
        "• Unstrukturierter Überblick: Struktur, Qualität, Ausreißer",
        "• Keine lokale, portable Lösung ohne Installation",
        "• Cloud-Abhängigkeit oft nicht gewünscht",
        "• Originaldatei muss geschützt sein",
        "→ Lösung: Schneller, lokaler, offline Überblick"
    ]
)

# 3. PRODUKTZIEL
add_content_slide(
    "Produktziel",
    [
        "✓ Schneller Überblick ohne Setup",
        "✓ Vollständig portabel und lokal",
        "✓ Originaldatei geschützt",
        "✓ Arbeitsfluss für Analysten: Laden → Prüfen → Bearbeiten → Export",
        "✓ Fokus auf Fachlichkeit, nicht Komplexität",
        "✓ Windows-nativer Desktop-Workflow"
    ]
)

# 4. FUNKTIONSUMFANG
add_content_slide(
    "Funktionsumfang",
    [
        "CSV-Verwaltung:",
        "  • CSV laden mit automatischer Encoding/Separator-Erkennung",
        "  • Beispieldatei (Schnellstart)",
        "",
        "Analyse & Vorschau:",
        "  • Strukturprofil und Datenqualität",
        "  • Kennzahlen und Diagramme",
        "  • Tabellenansicht mit Scrolling",
        "",
        "Arbeitstools:",
        "  • Mehrfach-Filter, Schnellansichten",
        "  • Spalten-Bearbeitung (Umbenennen, Löschen, Typ-Overrides)",
        "  • CSV-Export & HTML-Bericht (aktive Sicht)"
    ]
)

# 5. PRODUKTLOGIK
add_content_slide(
    "Produktlogik & Schutzkonzept",
    [
        "Originaldatei → Arbeitskopie",
        "  • Originaldatei wird beim Laden gespeichert",
        "  • Alle Änderungen nur auf Arbeitskopie",
        "",
        "Aktive Sicht",
        "  • Filter und Transformationen definieren aktive Sicht",
        "  • Kennzahlen, Diagramme, Export beziehen sich auf aktive Sicht",
        "",
        "Export & Bericht",
        "  • CSV-Export speichert aktive Sicht als neue Datei",
        "  • HTML-Bericht dokumentiert aktive Sicht",
        "",
        "Schutz",
        "  • Originaldatei wird nie verändert oder überschrieben"
    ]
)

# 6. OBERFLÄCHE / PRODUKTANSICHT
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_DARK_BG

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "Produktoberfläche"
for paragraph in title_frame.paragraphs:
    paragraph.font.size = Pt(40)
    paragraph.font.bold = True
    paragraph.font.color.rgb = COLOR_GREEN_ACCENT

line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
line.line.color.rgb = COLOR_GREEN_ACCENT
line.line.width = Pt(2)

# Screenshot
if RELEASE_SCREENSHOT.exists():
    try:
        slide.shapes.add_picture(str(RELEASE_SCREENSHOT), Inches(0.5), Inches(1.3), width=Inches(9))
    except Exception as e:
        # Fallback: Text
        text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = text_box.text_frame
        tf.text = "• Toolbar: CSV öffnen, Bericht exportieren, Info/Schnellstart\n• Tabs: Überblick, Kennzahlen, Diagramme, Tabelle, Bearbeiten, Export\n• Dunkles ModernDesign mit grünen Akzenten\n• Responsive Layouts für große und kleine Dateien"
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = COLOR_LIGHT_TEXT

# 7. TECHNIKBLOCK - Stack
add_content_slide(
    "Technik-Stack",
    [
        ("Programmiersprache", True),
        "  Python 3.10+",
        "",
        ("UI-Framework", True),
        "  PySide6 (Qt6 für Python)",
        "",
        ("Datenverwaltung", True),
        "  pandas (DataFrames, Analyse)",
        "  numpy (numerische Berechnung)",
        "",
        ("Visualisierung", True),
        "  matplotlib (Diagramme, Charts)",
        "",
        ("Qualitätssicherung", True),
        "  pytest (automatisierte Tests)"
    ]
)

# 8. TECHNIKBLOCK - Build & Release
add_content_slide(
    "Build & Release",
    [
        ("Packaging", True),
        "  PyInstaller One-Folder-Release",
        "  → MinAn_1_4/ Ordner (portable)",
        "",
        ("Release-Struktur", True),
        "  ├ MinAn.exe (Hauptanwendung)",
        "  ├ _internal/ (Python Runtime, Libraries)",
        "  ├ output/ (Berichte und CSV-Exporte)",
        "  └ _internal/sample_data/ (Beispieldatei)",
        "",
        ("Versionierung", True),
        "  Windows Version Info: 1.4.0.0",
        "  Productname: MinAn 1.4 - CSV-Schnellanalyse"
    ]
)

# 9. ARCHITEKTUR
add_two_column_slide(
    "Projektarchitektur",
    [
        ("Schichtmodell", True),
        "",
        "UI-Layer",
        "  main_window",
        "  dialogs, widgets",
        "  Qt-Models",
        "",
        "Services",
        "  import, profile",
        "  quality, chart",
        "  transform, export",
        "  report"
    ],
    [
        ("SessionState & Domain", True),
        "",
        "SessionState",
        "  original_df",
        "  working_df",
        "  filters & views",
        "",
        "Domain Models",
        "  ColumnProfile",
        "  DataQuality",
        "  enums.py"
    ]
)

# 10. QUALITÄTSSICHERUNG
add_content_slide(
    "Qualitätssicherung",
    [
        ("Automatisierte Tests", True),
        "  44+ Pytest-Tests abdecken:",
        "  • CSV-Import mit verschiedenen Encodings",
        "  • Profilierung (Zeilen, Spalten, Typen)",
        "  • Datenqualität (Missing, Dubletten)",
        "  • Transformationen (Filter, Markierungen)",
        "  • Export und Bericht-Generierung",
        "  • Originalschutz und Session-State",
        "",
        ("Manuelle Smoke-Tests", True),
        "  ✓ App startet und lädt Dateien",
        "  ✓ Alle Tabs funktionieren",
        "  ✓ Export und Bericht-Generierung"
    ]
)

# 11. RELEASE & PORTABILITÄT
add_content_slide(
    "Release & Portabilität",
    [
        ("One-Folder-Prinzip", True),
        "  dist/MinAn_1_4/ ist vollständig portabel",
        "  Keine Systeminstallation erforderlich",
        "",
        ("Ausgabeordner", True),
        "  output/reports/ → HTML-Berichte",
        "  output/csv/ → CSV-Exporte",
        "",
        ("Beispieldatei", True),
        "  test_csv_deutsch_200x15.csv",
        "  200 Zeilen, 15 Spalten, deutsch",
        "  → Schnellstart und Demo",
        "",
        ("Schnellstart", True),
        "  Info-Dialog mit Beispieldatei-Button"
    ]
)

# 12. PRODUKTREIFE & FEINSCHLIFF
add_content_slide(
    "Produktreife - Feinschliff",
    [
        ("UI/UX Politur", True),
        "  ✓ Dark Mode mit grünen Akzenten",
        "  ✓ Responsive Layouts",
        "  ✓ Klare, konsistente Toolbar",
        "  ✓ Info-Dialog und Schnellstart",
        "",
        ("Produktmetadaten", True),
        "  ✓ Produktname: MinAn 1.4 - CSV-Schnellanalyse",
        "  ✓ Company: MinAn Software",
        "  ✓ Icon: ModernGreen Analyse-Symbol",
        "",
        ("Fertigstellung 1.4", True),
        "  ✓ Alle Tests grün",
        "  ✓ Release-Build erfolgreich",
        "  ✓ Portable EXE lauffähig"
    ]
)

# 13. FAZIT
add_content_slide(
    "Fazit",
    [
        ("Was ist MinAn 1.4?", True),
        "Ein reifes, produktives Analysewerkzeug für CSV-Daten",
        "",
        ("Kernstärken", True),
        "  ✓ Lokal, portabel, sicher",
        "  ✓ Schneller Überblick ohne Setup",
        "  ✓ Originalschutz durch Arbeitskopien-Prinzip",
        "  ✓ Professioneller Export und Bericht",
        "  ✓ Solide Test-Basis und Release-Qualität",
        "",
        ("Projektstand", True),
        "  ✓ Version 1.4 ist Abschluss von Blocks 1–2",
        "  ✓ Produktionsreifer Stand",
        "  ✓ Bereit für Einsatz und Weiterentwicklung"
    ]
)

# SPEICHERN
output_path = PROJECT_ROOT / "presentation" / "MinAn_1_4_Abschlusspraesentation.pptx"
output_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(output_path))

print(f"[OK] Präsentation erstellt: {output_path}")
print(f"[OK] Folien: {len(prs.slides)}")
